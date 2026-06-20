"""Document ingestion for Course Content Studio.

Turn an uploaded file (PDF / PPTX / HTML / Markdown / text) into plain text plus
a list of *blocks* — small, source-located pieces of the document — so the rest
of the pipeline can chunk, embed, and cite back to where each fact came from.

The single entry point is :func:`extract`, which dispatches by file extension::

    out = extract("lecture.pdf", raw_bytes)
    out["text"]    # the whole document as one string
    out["blocks"]  # [{"text": "...", "loc": "p.1"}, ...]

``loc`` is a human-readable source locator that flows all the way through to the
citations on generated quiz items (``p.3``, ``slide 5``, ``html``, ``doc``).

Heavy parsers (``pdfplumber``, ``pypdf``, ``python-pptx``, ``beautifulsoup4``,
``markdownify``) are imported lazily *inside* the handlers so this module always
imports — and the page can show a friendly "install X" message instead of
crashing — even when an optional dependency is missing.
"""
from __future__ import annotations

import io
import os


def extract(filename: str, data: bytes) -> dict:
    """Extract text + located blocks from an uploaded file.

    Args:
        filename: the original file name (used only for its extension).
        data: the raw file bytes.

    Returns:
        ``{"text": str, "blocks": [{"text": str, "loc": str}, ...]}``. The
        ``blocks`` list may be empty for an empty or unreadable file, but this
        function never raises for empty/garbled content — only for an
        *unsupported* extension.

    Raises:
        ValueError: if the file extension is not one we know how to read.
    """
    ext = os.path.splitext(filename or "")[1].lower().lstrip(".")
    data = data or b""

    if ext == "pdf":
        blocks = _extract_pdf(data)
    elif ext == "pptx":
        blocks = _extract_pptx(data)
    elif ext in ("html", "htm"):
        blocks = _extract_html(data)
    elif ext in ("md", "markdown", "txt"):
        blocks = _extract_text(data)
    else:
        raise ValueError(
            f"Unsupported file type: {ext or '(none)'!r}. "
            "Upload a .pdf, .pptx, .html, .md, or .txt file."
        )

    text = "\n\n".join(b["text"] for b in blocks if b.get("text", "").strip())
    return {"text": text, "blocks": blocks}


# --- per-format handlers --------------------------------------------------

def _extract_pdf(data: bytes) -> list[dict]:
    """One block per page via pdfplumber, falling back to pypdf.

    Empty/garbled pages contribute nothing (no crash). If neither library is
    installed we raise a clear ImportError telling the caller what to add.
    """
    if not data:
        return []

    # Preferred: pdfplumber (better layout-aware text extraction).
    try:
        import pdfplumber

        blocks: list[dict] = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for n, page in enumerate(pdf.pages, start=1):
                try:
                    txt = page.extract_text() or ""
                except Exception:
                    txt = ""
                txt = txt.strip()
                if txt:
                    blocks.append({"text": txt, "loc": f"p.{n}"})
        return blocks
    except ImportError:
        pass  # try pypdf below
    except Exception:
        # pdfplumber installed but failed on this file — try pypdf as a fallback.
        pass

    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        blocks = []
        for n, page in enumerate(reader.pages, start=1):
            try:
                txt = page.extract_text() or ""
            except Exception:
                txt = ""
            txt = txt.strip()
            if txt:
                blocks.append({"text": txt, "loc": f"p.{n}"})
        return blocks
    except ImportError as e:
        raise ImportError(
            "Reading PDFs needs 'pdfplumber' (preferred) or 'pypdf'. "
            "Add one to requirements.txt."
        ) from e


def _extract_pptx(data: bytes) -> list[dict]:
    """One block per slide: every shape's text plus the slide's speaker notes."""
    if not data:
        return []
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError(
            "Reading .pptx needs 'python-pptx'. Add it to requirements.txt."
        ) from e

    prs = Presentation(io.BytesIO(data))
    blocks: list[dict] = []
    for n, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        parts.append(t)
            except Exception:
                continue
        # Speaker notes are a separate, very useful source of context.
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    parts.append(f"[notes] {notes}")
        except Exception:
            pass
        txt = "\n".join(parts).strip()
        if txt:
            blocks.append({"text": txt, "loc": f"slide {n}"})
    return blocks


def _extract_html(data: bytes) -> list[dict]:
    """Whole-document text via BeautifulSoup (markdownify if available).

    Scripts/styles are dropped; one block with ``loc="html"`` is returned.
    """
    if not data:
        return []
    raw = _decode(data)
    try:
        from bs4 import BeautifulSoup
    except ImportError as e:
        raise ImportError(
            "Reading HTML needs 'beautifulsoup4'. Add it to requirements.txt."
        ) from e

    soup = BeautifulSoup(raw, "html.parser")
    for bad in soup(["script", "style", "noscript"]):
        bad.decompose()

    text = ""
    # markdownify keeps headings/lists as structure, which reads better; if it
    # isn't installed we fall back to plain get_text().
    try:
        from markdownify import markdownify as md

        text = md(str(soup)).strip()
    except Exception:
        text = soup.get_text(separator="\n").strip()

    # Collapse runs of blank lines so chunks aren't padded with whitespace.
    text = "\n".join(line.rstrip() for line in text.splitlines())
    text = "\n".join(s for s in text.split("\n") if s.strip() or True).strip()
    if not text:
        return []
    return [{"text": text, "loc": "html"}]


def _extract_text(data: bytes) -> list[dict]:
    """Plain text / Markdown — decode UTF-8 (lenient) into a single block."""
    if not data:
        return []
    text = _decode(data).strip()
    if not text:
        return []
    return [{"text": text, "loc": "doc"}]


def _decode(data: bytes) -> str:
    """Best-effort UTF-8 decode that never raises on garbled bytes."""
    try:
        return data.decode("utf-8")
    except Exception:
        return data.decode("utf-8", errors="replace")
